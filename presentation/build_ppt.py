"""
MV Kobra AI · Generador de presentación gerencial (PPTX)
==================================================
Crea un deck comercial atractivo para vender MV Kobra AI a la gerencia de cualquier
empresa con cartera de cobranzas. Toma los KPIs reales del pipeline
(outputs/kobra_bundle.json) y capturas del dashboard (assets/).

Uso:
    python presentation/build_ppt.py
"""
import json
import os
import sys

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

ROOT = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
ASSETS = os.path.join(ROOT, "assets")
OUT = os.path.join(ROOT, "presentation", "MVKobraAI_Presentacion_Gerencial.pptx")

# Paleta de marca — desde marketing/marca.py, que es la fuente única.
# Antes estaba escrita a mano acá y había divergido de la del sitio: la
# presentación gerencial y la landing no parecían de la misma empresa. El
# violeta y el amarillo que se usaban para distinguir categorías pasaron a ser
# el azul y el ámbar de la marca, que cumplen el mismo rol y sí aparecen en el
# resto del material.
sys.path.insert(0, ROOT)
from marketing import marca as _marca  # noqa: E402


def _c(nombre):
    return RGBColor(*_marca.rgb(_marca.MARCA[nombre]))


DARK = _c("navy")
PANEL = _c("panel")
GREEN = _c("green")
PURPLE = _c("blue")      # acento 2 (era violeta #6C5CE7, ajeno a la marca)
WHITE = _c("ink")
GREY = _c("muted")
YELLOW = _c("amber")     # acento 3 (era amarillo #FDCB6E, ídem)

W, H = Inches(13.333), Inches(7.5)


def _load_kpis():
    p = os.path.join(ROOT, "outputs", "kobra_bundle.json")
    if os.path.exists(p):
        with open(p, encoding="utf-8") as f:
            return json.load(f)
    return {"kpis": {}, "metrics": {}}


def _bg(slide, color=DARK):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = color


def _box(slide, x, y, w, h, fill=None, line=None):
    from pptx.enum.shapes import MSO_SHAPE
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    if fill is None:
        shp.fill.background()
    else:
        shp.fill.solid(); shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line; shp.line.width = Pt(1)
    shp.shadow.inherit = False
    return shp


def _text(slide, x, y, w, h, text, size=18, color=WHITE, bold=False,
          align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, font="Segoe UI"):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    lines = text.split("\n")
    for i, ln in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        r = p.add_run(); r.text = ln
        r.font.size = Pt(size); r.font.bold = bold
        r.font.color.rgb = color; r.font.name = font
    return tb


def _bullets(slide, x, y, w, h, items, size=17, color=WHITE, gap=6):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame; tf.word_wrap = True
    for i, (txt, c) in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(gap)
        r = p.add_run(); r.text = "▸  " + txt
        r.font.size = Pt(size); r.font.color.rgb = c or color; r.font.name = "Segoe UI"
    return tb


def _fmt_m(v):
    return f"$U {v/1e6:,.1f}M" if v else "—"


def build():
    data = _load_kpis()
    k = data.get("kpis", {})
    m = data.get("metrics", {})
    prs = Presentation()
    prs.slide_width = W; prs.slide_height = H
    blank = prs.slide_layouts[6]

    # ---------- 1 · Portada -------------------------------------------------
    s = prs.slides.add_slide(blank); _bg(s)
    _box(s, 0, Inches(3.15), W, Inches(1.2), fill=None)
    logo = os.path.join(ROOT, "assets", "brand", "mv_icon.png")
    if os.path.exists(logo):
        s.shapes.add_picture(logo, Inches(5.92), Inches(0.85), height=Inches(1.5))
    _text(s, Inches(1), Inches(2.4), Inches(11.3), Inches(1.2),
          "MV Kobra AI", size=66, color=GREEN, bold=True, align=PP_ALIGN.CENTER)
    _text(s, Inches(1), Inches(3.7), Inches(11.3), Inches(0.7),
          "Plataforma de Cobranzas Inteligentes", size=28, color=WHITE,
          align=PP_ALIGN.CENTER)
    _text(s, Inches(1), Inches(4.5), Inches(11.3), Inches(0.6),
          "ProbPago · Agente IA Negociador · Priorización por valor esperado",
          size=17, color=GREY, align=PP_ALIGN.CENTER)
    _text(s, Inches(1), Inches(6.6), Inches(11.3), Inches(0.5),
          "Propuesta gerencial · Uruguay · Demo con datos sintéticos (sin datos de clientes)",
          size=12, color=GREY, align=PP_ALIGN.CENTER)

    # ---------- 2 · El problema --------------------------------------------
    s = prs.slides.add_slide(blank); _bg(s)
    _text(s, Inches(0.7), Inches(0.5), Inches(12), Inches(0.8),
          "El problema de la cobranza tradicional", size=32, color=WHITE, bold=True)
    _box(s, Inches(0.7), Inches(1.4), Inches(0.15), Inches(0.7), fill=GREEN)
    problems = [
        ("Se gestiona por igual a todos los deudores → se malgasta tiempo en casos irrecuperables.", GREY),
        ("Las decisiones de descuento y planes se toman \"a ojo\", sin datos.", GREY),
        ("No se sabe a quién llamar primero ni por qué canal.", GREY),
        ("Los guiones de negociación no están estandarizados ni optimizados.", GREY),
        ("El reporting gerencial es lento, manual y poco accionable.", GREY),
    ]
    _bullets(s, Inches(0.9), Inches(2.0), Inches(11.5), Inches(4.5), problems, size=20, gap=14)
    _text(s, Inches(0.9), Inches(6.3), Inches(11.5), Inches(0.8),
          "Resultado: menor recupero, más costo operativo y peor experiencia del cliente.",
          size=18, color=YELLOW, bold=True)

    # ---------- 3 · La solución (3 pilares) --------------------------------
    s = prs.slides.add_slide(blank); _bg(s)
    _text(s, Inches(0.7), Inches(0.5), Inches(12), Inches(0.8),
          "La solución: MV Kobra AI", size=32, color=WHITE, bold=True)
    _box(s, Inches(0.7), Inches(1.4), Inches(0.15), Inches(0.7), fill=GREEN)
    cards = [
        ("🧠 ProbPago", GREEN,
         "Modelo de Machine Learning que estima la probabilidad de pago de cada deudor. "
         "Selección de modelos con validación cruzada + calibración.",
         "Se entrena y valida con la cartera real del cliente"),
        ("🤖 Agente + Copiloto en Vivo", PURPLE,
         "Recomienda estrategia, descuento, canal y guion. En vivo: analiza el sentimiento del "
         "cliente (voz/WhatsApp) y asesora al gestor en tiempo real.",
         "Maximiza recupero minimizando la quita"),
        ("📊 Dashboard Gerencial", YELLOW,
         "KPIs, filtros, gráficos y priorización de cartera por valor esperado de recupero.",
         "Exportable a Excel / CSV para reporting"),
    ]
    cw = Inches(3.9); gap = Inches(0.25); x0 = Inches(0.7); y0 = Inches(2.1)
    for i, (title, col, desc, foot) in enumerate(cards):
        x = x0 + i * (cw + gap)
        _box(s, x, y0, cw, Inches(4.2), fill=PANEL, line=col)
        _text(s, x + Inches(0.3), y0 + Inches(0.35), cw - Inches(0.6), Inches(0.8),
              title, size=22, color=col, bold=True)
        _text(s, x + Inches(0.3), y0 + Inches(1.4), cw - Inches(0.6), Inches(2),
              desc, size=16, color=WHITE)
        _text(s, x + Inches(0.3), y0 + Inches(3.4), cw - Inches(0.6), Inches(0.7),
              foot, size=13, color=col, bold=True)

    # ---------- 4 · Cómo funciona (flujo) ----------------------------------
    s = prs.slides.add_slide(blank); _bg(s)
    _text(s, Inches(0.7), Inches(0.5), Inches(12), Inches(0.8),
          "Cómo funciona · flujo end-to-end", size=32, color=WHITE, bold=True)
    _box(s, Inches(0.7), Inches(1.4), Inches(0.15), Inches(0.7), fill=GREEN)
    steps = [
        ("1 · Cartera", "Se carga la base de deudores (o se conecta al core).", GREEN),
        ("2 · ProbPago", "El modelo asigna probabilidad de pago a cada caso.", PURPLE),
        ("3 · Negociador IA", "Define estrategia, descuento, canal y guion.", YELLOW),
        ("4 · Priorización", "Ordena por valor esperado de recupero (UYU).", GREEN),
        ("5 · Acción + Reporting", "Gestores ejecutan; gerencia mide en el dashboard.", PURPLE),
    ]
    y = Inches(2.1)
    for title, desc, col in steps:
        _box(s, Inches(0.9), y, Inches(3.2), Inches(0.82), fill=PANEL, line=col)
        _text(s, Inches(1.1), y + Inches(0.12), Inches(3), Inches(0.6),
              title, size=17, color=col, bold=True)
        _text(s, Inches(4.4), y + Inches(0.12), Inches(8), Inches(0.6),
              desc, size=17, color=WHITE)
        y += Inches(0.98)

    # ---------- 5 · Resultados / KPIs --------------------------------------
    s = prs.slides.add_slide(blank); _bg(s)
    _text(s, Inches(0.7), Inches(0.5), Inches(12), Inches(0.8),
          "Impacto sobre la cartera (demo)", size=32, color=WHITE, bold=True)
    _box(s, Inches(0.7), Inches(1.4), Inches(0.15), Inches(0.7), fill=GREEN)
    kpis = [
        (f"{k.get('deudores',0):,}", "Deudores analizados", GREEN),
        (_fmt_m(k.get('cartera_total_uyu',0)), "Cartera total", WHITE),
        (_fmt_m(k.get('recupero_esperado_uyu',0)),
         f"Recupero esperado ({k.get('tasa_recupero_esperada',0):.0%})", GREEN),
        (f"{k.get('probpago_promedio',0):.0%}", "ProbPago promedio", PURPLE),
        (_fmt_m(k.get('en_riesgo_uyu',0)), "Cartera en riesgo (baja prop.)", YELLOW),
        (f"{m.get('lift_decil10','1.7')}x", "Lift del top decil vs. base", GREEN),
    ]
    cw = Inches(3.9); ch = Inches(1.9); gx = Inches(0.25); gy = Inches(0.3)
    x0 = Inches(0.7); y0 = Inches(2.1)
    for i, (val, lbl, col) in enumerate(kpis):
        r, c = divmod(i, 3)
        x = x0 + c * (cw + gx); yy = y0 + r * (ch + gy)
        _box(s, x, yy, cw, ch, fill=PANEL, line=None)
        _text(s, x, yy + Inches(0.3), cw, Inches(0.9), val, size=34,
              color=col, bold=True, align=PP_ALIGN.CENTER)
        _text(s, x, yy + Inches(1.25), cw, Inches(0.5), lbl, size=14,
              color=GREY, align=PP_ALIGN.CENTER)
    _text(s, Inches(0.7), Inches(6.7), Inches(12), Inches(0.5),
          "Datos sintéticos ilustrativos: muestran cómo MV Kobra AI prioriza el esfuerzo donde está "
          "el valor recuperable. Con la cartera real, estos KPIs se calculan de verdad.",
          size=14, color=GREY, align=PP_ALIGN.CENTER)

    # ---------- 6/7/8 · Capturas del dashboard -----------------------------
    shots = [
        ("dashboard_overview.png", "Dashboard · Visión general",
         "KPIs en vivo, filtros dinámicos y gráficos de cartera vs. recupero."),
        ("dashboard_negociador.png", "Agente IA Negociador",
         "Estrategias recomendadas y simulador de negociación con guion por deudor."),
        ("dashboard_cartera.png", "Cartera priorizada + Export",
         "Ranking por valor esperado y exportación a Excel/CSV para reporting."),
        ("dashboard_copiloto_detalle.png", "Copiloto de Negociación en Vivo",
         "Sentimiento del cliente turno a turno, técnicas del gestor y la próxima jugada sugerida."),
        ("realtime_copiloto.png", "Asistencia en vivo durante la llamada",
         "Escucha la llamada telefónica, transcribe y asesora al gestor en tiempo real."),
        ("dashboard_voz2.png", "Voz: diarización + emoción acústica",
         "Separa quién habla y detecta la tensión del cliente en la voz, más allá de las palabras."),
        ("realtime_telefonia.png", "Integración telefónica + transcripción por hablante",
         "Conecta al audio de la central (Avaya, Genesys…); transcribe por hablante y fusiona voz + texto."),
        ("realtime_stream.png", "Streaming en vivo durante la llamada",
         "SIPREC / Twilio Media Streams / Avaya DMCC: asesoría turno a turno mientras la llamada ocurre."),
        ("dashboard_gestores.png", "Gestores & Evolución",
         "Impacto de MV Kobra AI en calidad, conversión y recupero; evolución mes a mes por gestor."),
    ]
    for img, title, sub in shots:
        s = prs.slides.add_slide(blank); _bg(s)
        _text(s, Inches(0.7), Inches(0.35), Inches(12), Inches(0.7),
              title, size=28, color=GREEN, bold=True)
        _text(s, Inches(0.7), Inches(1.0), Inches(12), Inches(0.5),
              sub, size=15, color=GREY)
        path = os.path.join(ASSETS, img)
        if os.path.exists(path):
            _box(s, Inches(0.55), Inches(1.55), Inches(12.25), Inches(5.55),
                 fill=None, line=PURPLE)
            s.shapes.add_picture(path, Inches(0.7), Inches(1.7), width=Inches(11.95))

    # ---------- 8b · Impacto de MV Kobra AI en la gestión ------------------------
    imp_path = os.path.join(ROOT, "outputs", "impacto_kobra.json")
    if os.path.exists(imp_path):
        with open(imp_path, encoding="utf-8") as fh:
            ik = json.load(fh)
        s = prs.slides.add_slide(blank); _bg(s)
        _text(s, Inches(0.7), Inches(0.5), Inches(12), Inches(0.8),
              "Cómo se mide el impacto de MV Kobra AI", size=32, color=WHITE, bold=True)
        _box(s, Inches(0.7), Inches(1.4), Inches(0.15), Inches(0.7), fill=GREEN)
        _text(s, Inches(0.9), Inches(1.7), Inches(11.5), Inches(0.6),
              "Metodología: gestores con MV Kobra AI vs. grupo de control, mismo período. "
              "Cifras ILUSTRATIVAS sobre datos sintéticos — no son resultados medidos.",
              size=15, color=YELLOW)
        cards = [
            (f"+{ik['uplift_calidad']:.1f}", "puntos de calidad de gestión", GREEN),
            (f"+{ik['uplift_conversion']*100:.1f} pp", "tasa de conversión", PURPLE),
            (f"+{ik['uplift_recupero']*100:.1f} pp", "tasa de recupero", YELLOW),
        ]
        cw = Inches(3.9); gap = Inches(0.25); x0 = Inches(0.7); y0 = Inches(2.6)
        for i, (val, lbl, col) in enumerate(cards):
            x = x0 + i * (cw + gap)
            _box(s, x, y0, cw, Inches(2.4), fill=PANEL, line=col)
            _text(s, x, y0 + Inches(0.55), cw, Inches(1), val, size=44,
                  color=col, bold=True, align=PP_ALIGN.CENTER)
            _text(s, x, y0 + Inches(1.7), cw, Inches(0.5), lbl, size=15,
                  color=GREY, align=PP_ALIGN.CENTER)
        _text(s, Inches(0.9), Inches(5.6), Inches(11.5), Inches(1),
              f"Con MV Kobra AI: calidad {ik['con_kobra']['calidad_prom']:.0f} · "
              f"conversión {ik['con_kobra']['tasa_conversion']:.0%} · "
              f"recupero {ik['con_kobra']['tasa_recupero']:.0%}   |   "
              f"Sin MV Kobra AI: calidad {ik['sin_kobra']['calidad_prom']:.0f} · "
              f"conversión {ik['sin_kobra']['tasa_conversion']:.0%} · "
              f"recupero {ik['sin_kobra']['tasa_recupero']:.0%}",
              size=14, color=WHITE, align=PP_ALIGN.CENTER)
        _text(s, Inches(0.9), Inches(6.4), Inches(11.5), Inches(0.6),
              "En un piloto real, este mismo análisis (control vs. adopción) produce el impacto "
              "MEDIDO — esa es la evidencia que se presenta, no esta simulación.",
              size=14, color=GREEN, align=PP_ALIGN.CENTER)

    # ---------- 9 · Por qué MV Kobra AI / diferenciales --------------------------
    s = prs.slides.add_slide(blank); _bg(s)
    _text(s, Inches(0.7), Inches(0.5), Inches(12), Inches(0.8),
          "Por qué MV Kobra AI", size=32, color=WHITE, bold=True)
    _box(s, Inches(0.7), Inches(1.4), Inches(0.15), Inches(0.7), fill=GREEN)
    why = [
        ("Más recupero: prioriza por valor esperado, no por antigüedad de la deuda.", GREEN),
        ("Asistencia en vivo: análisis de sentimiento de voz/WhatsApp que guía al gestor en la negociación.", GREEN),
        ("Menos costo: automatiza estrategia y guiones; el gestor solo ejecuta.", GREEN),
        ("Decisiones con datos: ML de probabilidad de pago con selección de modelos y calibración.", GREEN),
        ("Listo para Uruguay: pensado para el mercado local, sin datos personales sensibles.", GREEN),
        ("Agnóstico de industria: banca, financieras, retail, telco, utilities, fintech.", GREEN),
    ]
    _bullets(s, Inches(0.9), Inches(2.0), Inches(11.5), Inches(4.5), why, size=19, gap=14)

    # ---------- 9b · Listo para producción ---------------------------------
    s = prs.slides.add_slide(blank); _bg(s)
    _text(s, Inches(0.7), Inches(0.5), Inches(12), Inches(0.8),
          "Listo para producción", size=32, color=WHITE, bold=True)
    _box(s, Inches(0.7), Inches(1.4), Inches(0.15), Inches(0.7), fill=GREEN)
    prod = [
        ("🐳 Despliegue en 1 comando", "docker compose up: dashboard + servicio de audio en vivo, con datos y modelo autogenerados."),
        ("🔑 API keys persistentes", "Se cargan una vez desde la pestaña Configuración (o por variables de entorno) y quedan guardadas."),
        ("☎️ Integración telefónica", "Conectores para Avaya/Genesys/Twilio (SIPREC · Media Streams · DMCC) — grabación o streaming en vivo."),
        ("✅ Calidad asegurada", "Suite de tests + CI (GitHub Actions) y reentrenamiento programado del modelo ProbPago."),
        ("🔒 Sin datos sensibles", "Dataset sintético, sin nombres de clientes. Reemplazable por la cartera real de la empresa."),
    ]
    y = Inches(2.0)
    for titulo, desc, in [(p[0], p[1]) for p in prod]:
        _box(s, Inches(0.9), y, Inches(4.6), Inches(0.86), fill=PANEL, line=GREEN)
        _text(s, Inches(1.1), y + Inches(0.16), Inches(4.3), Inches(0.6),
              titulo, size=15, color=GREEN, bold=True)
        _text(s, Inches(5.8), y + Inches(0.1), Inches(6.8), Inches(0.7),
              desc, size=14, color=WHITE)
        y += Inches(0.98)

    # ---------- 10 · Cierre -------------------------------------------------
    s = prs.slides.add_slide(blank); _bg(s)
    if os.path.exists(logo):
        s.shapes.add_picture(logo, Inches(6.17), Inches(1.15), height=Inches(1.0))
    _text(s, Inches(1), Inches(2.5), Inches(11.3), Inches(1),
          "Convierta su cartera en recupero.", size=40, color=WHITE, bold=True,
          align=PP_ALIGN.CENTER)
    _text(s, Inches(1), Inches(3.7), Inches(11.3), Inches(0.8),
          "MV Kobra AI · Cobranzas Inteligentes", size=26, color=GREEN, bold=True,
          align=PP_ALIGN.CENTER)
    _text(s, Inches(1), Inches(4.6), Inches(11.3), Inches(0.6),
          "Demo lista para ejecutar · Solicite una prueba con su propia cartera.",
          size=16, color=GREY, align=PP_ALIGN.CENTER)

    prs.save(OUT)
    print(f"[OK] Presentación generada: {OUT}  ({len(prs.slides.__iter__.__self__._sldIdLst)} slides)")


if __name__ == "__main__":
    build()
